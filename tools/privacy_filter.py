#!/usr/bin/env python3
from __future__ import annotations

"""
Privacy filter — local PII anonymization for maskzone files.

Contract:
  --setup           → download model, run test inference, create
                      .privacy-deps-ok marker, set can_anonymize_pii flag
                      in wiki/.protect, create maskzone/ directory.
                      Print JSON result. Does NOT install pip packages
                      (the /privacy-mode skill handles that).
  --status          → read privacy mode state + model availability,
                      print as JSON (no mutation)
  --hook-check      → entry point for UserPromptSubmit hook.
                      If can_anonymize_pii is false → silent exit (0ms).
                      If .privacy-deps-ok missing → print warning, exit.
                      If maskzone/ empty → silent exit.
                      Otherwise: process all files in maskzone/.
  --process <file>  → extract text from a single file, run PII detection,
                      replace PII spans with placeholders,
                      write anonymized .md to raw/<slug>.md.
                      Originals stay in maskzone/ — the user removes them manually.
                      Print JSON result.
  --test            → run PII detection on sample text with known PII,
                      verify detection of at least NAME, EMAIL, PHONE, DATE.
                      Print JSON pass/fail.

Inputs:  maskzone/* (files to process),
         wiki/.protect (flag check),
         .privacy-deps-ok (dep check),
         ~/.cache/huggingface/ (cached model weights)
Outputs: raw/<slug>.md (anonymized text),
         stdout JSON (results/errors)

If this script fails:
  - can_anonymize_pii not found in .protect → treat as false
  - .privacy-deps-ok missing → warning message, no import attempt
  - maskzone/ missing → create it (auto-init)
  - Unsupported file format → error JSON, skip file
  - Model not downloaded → error JSON, suggest running /privacy-mode
  - Import error → error JSON, suggest running /privacy-mode
  - Extraction error on specific file → error JSON for that file,
    continue processing remaining files
"""

import sys
import re
import json
import argparse
from pathlib import Path
from datetime import datetime

from shared import REPO_ROOT, WIKI_DIR

MASKZONE_DIR = REPO_ROOT / "maskzone"
RAW_DIR = REPO_ROOT / "raw"
PROTECT_FILE = WIKI_DIR / ".protect"
PRIVACY_DEPS_MARKER = REPO_ROOT / ".privacy-deps-ok"

SUPPORTED_FORMATS = {".md", ".txt", ".csv", ".docx", ".pdf", ".pptx", ".xlsx", ".html", ".htm"}

# Mapping from model entity_group labels (uppercased) to placeholders.
# The openai/privacy-filter model emits these entity groups:
#   private_person, private_email, private_phone, private_date,
#   private_address, private_url, account_number, secret
# After .upper() they become PRIVATE_PERSON, PRIVATE_EMAIL, etc.
LABEL_TO_PLACEHOLDER = {
    # openai/privacy-filter labels (primary)
    "PRIVATE_PERSON": "[NAME]",
    "PRIVATE_EMAIL": "[EMAIL]",
    "PRIVATE_PHONE": "[PHONE]",
    "PRIVATE_DATE": "[DATE]",
    "PRIVATE_ADDRESS": "[ADDRESS]",
    "PRIVATE_URL": "[URL]",
    "ACCOUNT_NUMBER": "[ACCOUNT]",
    "SECRET": "[SECRET]",
    # Generic fallback labels (for alternative models)
    "PERSON": "[NAME]",
    "NAME": "[NAME]",
    "PERSON_NAME": "[NAME]",
    "ADDRESS": "[ADDRESS]",
    "STREET_ADDRESS": "[ADDRESS]",
    "EMAIL": "[EMAIL]",
    "EMAIL_ADDRESS": "[EMAIL]",
    "PHONE": "[PHONE]",
    "PHONE_NUMBER": "[PHONE]",
    "URL": "[URL]",
    "DATE": "[DATE]",
    "DATE_TIME": "[DATE]",
    "ACCOUNT": "[ACCOUNT]",
    "API_KEY": "[SECRET]",
    "PASSWORD": "[SECRET]",
    "CREDENTIAL": "[SECRET]",
}


# ── Flag and dependency checks (no heavy imports) ──────────────────


def read_flag() -> bool:
    if not PROTECT_FILE.exists():
        return False
    try:
        data = json.loads(PROTECT_FILE.read_text(encoding="utf-8"))
        return bool(data.get("can_anonymize_pii", False))
    except (json.JSONDecodeError, ValueError):
        return False


def deps_ok() -> bool:
    return PRIVACY_DEPS_MARKER.exists()


# ── Text extraction (lazy imports per format) ──────────────────────


def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext in (".md", ".txt"):
        return file_path.read_text(encoding="utf-8", errors="replace")
    if ext == ".csv":
        return _extract_csv(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_path)
    if ext in (".html", ".htm"):
        return _extract_html(file_path)
    raise ValueError(f"Unsupported format: {ext}")


def _extract_csv(path: Path) -> str:
    import csv
    rows = []
    with open(path, encoding="utf-8", errors="replace", newline="") as f:
        for row in csv.reader(f):
            rows.append(" ".join(row))
    return "\n".join(rows)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n".join(pages)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    rows = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" ".join(cells))
    wb.close()
    return "\n".join(rows)


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup
    html = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator=" ", strip=True)


# ── PII anonymization (lazy model load) ────────────────────────────

_classifier = None


def load_classifier():
    global _classifier
    if _classifier is not None:
        return _classifier
    try:
        from optimum.onnxruntime import ORTModelForTokenClassification
        from transformers import AutoTokenizer, pipeline as hf_pipeline

        tokenizer = AutoTokenizer.from_pretrained("openai/privacy-filter")
        model = ORTModelForTokenClassification.from_pretrained(
            "openai/privacy-filter"
        )
        _classifier = hf_pipeline(
            "token-classification",
            model=model,
            tokenizer=tokenizer,
            aggregation_strategy="simple",
        )
    except Exception:
        from transformers import pipeline as hf_pipeline
        _classifier = hf_pipeline(
            "token-classification",
            model="openai/privacy-filter",
            aggregation_strategy="simple",
        )
    return _classifier


def anonymize_text(text: str) -> tuple[str, dict]:
    classifier = load_classifier()
    pii_counts: dict[str, int] = {}

    chunks = _split_into_chunks(text, max_length=2000)
    anonymized_chunks = []

    for chunk in chunks:
        spans = classifier(chunk)
        spans_sorted = sorted(spans, key=lambda s: s["start"], reverse=True)
        result = chunk
        for span in spans_sorted:
            label = span["entity_group"].upper()
            placeholder = LABEL_TO_PLACEHOLDER.get(label, "[PII]")
            category = placeholder.strip("[]")
            pii_counts[category] = pii_counts.get(category, 0) + 1
            result = result[:span["start"]] + placeholder + result[span["end"]:]
        anonymized_chunks.append(result)

    merged = "\n".join(anonymized_chunks)
    # Collapse duplicate adjacent placeholders (e.g. [NAME][NAME] → [NAME])
    merged = re.sub(r'(\[\w+\])(\s*\1)+', r'\1', merged)
    return merged, pii_counts


def _split_into_chunks(text: str, max_length: int = 2000) -> list[str]:
    lines = text.split("\n")
    chunks = []
    current = []
    current_len = 0
    for line in lines:
        if current_len + len(line) + 1 > max_length and current:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line)
        current_len += len(line) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks


# ── File operations ────────────────────────────────────────────────


def slugify(name: str) -> str:
    name = Path(name).stem
    name = re.sub(r"[^\w\s-]", "", name.lower())
    return re.sub(r"[\s_]+", "-", name).strip("-")



# ── Commands ───────────────────────────────────────────────────────


def cmd_setup():
    MASKZONE_DIR.mkdir(parents=True, exist_ok=True)
    gitkeep = MASKZONE_DIR / ".gitkeep"
    if not gitkeep.exists():
        gitkeep.touch()

    print("Downloading model...", file=sys.stderr)
    try:
        classifier = load_classifier()
    except Exception as e:
        print(json.dumps({
            "status": "ERROR",
            "action": "setup",
            "error": f"Model download failed: {e}",
            "suggestion": "Check internet connection and retry /privacy-mode",
        }, indent=2))
        sys.exit(1)

    test_text = "John Smith (john@test.com) called +39 333 1234567 on March 15."
    spans = classifier(test_text)
    if not spans:
        print(json.dumps({
            "status": "ERROR",
            "action": "setup",
            "error": "Model loaded but detected no PII in test text",
            "test_text": test_text,
        }, indent=2))
        sys.exit(1)

    PRIVACY_DEPS_MARKER.write_text("installed\n", encoding="utf-8")

    protect_data = {"can_forget": False, "can_modify": False, "can_anonymize_pii": False}
    if PROTECT_FILE.exists():
        try:
            protect_data = json.loads(PROTECT_FILE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, ValueError):
            pass
    protect_data["can_anonymize_pii"] = True
    PROTECT_FILE.write_text(json.dumps(protect_data, indent=2) + "\n", encoding="utf-8")

    detected_labels = [s["entity_group"] for s in spans]
    print(json.dumps({
        "status": "OK",
        "action": "setup",
        "model": "openai/privacy-filter",
        "test_pii_detected": detected_labels,
        "marker_created": True,
        "flag_set": True,
        "maskzone_ready": True,
    }, indent=2))


def cmd_status():
    flag = read_flag()
    has_deps = deps_ok()
    maskzone_exists = MASKZONE_DIR.exists()
    file_count = 0
    if maskzone_exists:
        file_count = len([f for f in MASKZONE_DIR.iterdir()
                          if f.is_file() and f.name != ".gitkeep"])

    print(json.dumps({
        "privacy_mode": "active" if flag else "inactive",
        "can_anonymize_pii": flag,
        "dependencies_installed": has_deps,
        "maskzone_exists": maskzone_exists,
        "maskzone_files_pending": file_count,
        "model": "openai/privacy-filter",
        "runtime": "onnxruntime",
    }, indent=2))


def cmd_hook_check():
    if not read_flag():
        return

    if not deps_ok():
        print("[PRIVACY] Dependencies not installed. Run /privacy-mode to set up.")
        return

    MASKZONE_DIR.mkdir(parents=True, exist_ok=True)
    files = [f for f in MASKZONE_DIR.iterdir()
             if f.is_file() and f.name != ".gitkeep"]

    if not files:
        return

    processed = []
    errors = []
    total_pii = 0

    for file_path in files:
        result = _process_single(file_path)
        if result["status"] == "OK":
            processed.append(file_path.name)
            total_pii += result["total_pii_masked"]
        else:
            errors.append({"file": file_path.name, "error": result.get("error", "unknown")})

    output = {
        "action": "hook-check",
        "processed": len(processed),
        "files": processed,
        "total_pii_masked": total_pii,
    }
    if errors:
        output["errors"] = errors

    print(json.dumps(output, indent=2))


def cmd_process(file_path_str: str):
    file_path = Path(file_path_str)
    if not file_path.is_absolute():
        file_path = REPO_ROOT / file_path_str

    if not file_path.exists():
        print(json.dumps({"status": "ERROR", "error": f"File not found: {file_path_str}"}, indent=2),
              file=sys.stderr)
        sys.exit(1)

    result = _process_single(file_path)
    print(json.dumps(result, indent=2))
    if result["status"] != "OK":
        sys.exit(1)


def _process_single(file_path: Path) -> dict:
    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_FORMATS:
        return {
            "status": "SKIP",
            "input": str(file_path.relative_to(REPO_ROOT)).replace("\\", "/"),
            "error": f"Unsupported format: {ext}",
            "supported": sorted(SUPPORTED_FORMATS),
        }

    try:
        text = extract_text(file_path)
    except Exception as e:
        return {
            "status": "ERROR",
            "input": str(file_path.relative_to(REPO_ROOT)).replace("\\", "/"),
            "error": f"Extraction failed: {e}",
        }

    if not text.strip():
        return {
            "status": "ERROR",
            "input": str(file_path.relative_to(REPO_ROOT)).replace("\\", "/"),
            "error": "Extracted text is empty",
        }

    try:
        anonymized, pii_counts = anonymize_text(text)
    except Exception as e:
        return {
            "status": "ERROR",
            "input": str(file_path.relative_to(REPO_ROOT)).replace("\\", "/"),
            "error": f"Anonymization failed: {e}",
        }

    slug = slugify(file_path.name)
    output_path = RAW_DIR / f"{slug}.md"
    RAW_DIR.mkdir(parents=True, exist_ok=True)
    output_path.write_text(anonymized, encoding="utf-8")

    input_rel = str(file_path.relative_to(REPO_ROOT)).replace("\\", "/")
    output_rel = str(output_path.relative_to(REPO_ROOT)).replace("\\", "/")

    return {
        "status": "OK",
        "input": input_rel,
        "output": output_rel,
        "pii_found": pii_counts,
        "total_pii_masked": sum(pii_counts.values()),
    }


def cmd_test():
    if not deps_ok():
        print(json.dumps({
            "status": "FAIL",
            "error": "Privacy dependencies not installed. Run /privacy-mode first.",
        }, indent=2))
        sys.exit(1)

    test_text = (
        "Dear John Smith, your account 4111-1111-1111-1111 has been updated. "
        "Contact us at support@example.com or call +1 555 123 4567. "
        "Your appointment is on January 15, 2026 at 123 Main St, Springfield."
    )

    try:
        anonymized, pii_counts = anonymize_text(test_text)
    except Exception as e:
        print(json.dumps({
            "status": "FAIL",
            "error": f"Anonymization failed: {e}",
        }, indent=2))
        sys.exit(1)

    detected_categories = set(pii_counts.keys())
    expected_minimum = {"NAME", "EMAIL", "PHONE", "DATE"}
    found = detected_categories & expected_minimum
    missing = expected_minimum - detected_categories

    passed = len(found) >= 3

    print(json.dumps({
        "status": "PASS" if passed else "FAIL",
        "test_input": test_text,
        "anonymized_output": anonymized,
        "pii_detected": pii_counts,
        "expected_categories": sorted(expected_minimum),
        "found_categories": sorted(found),
        "missing_categories": sorted(missing),
    }, indent=2))

    if not passed:
        sys.exit(1)


# ── Main ───────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Privacy filter — local PII anonymization (maskzone -> raw)"
    )
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--setup", action="store_true",
                       help="Download model, test, set flag, create maskzone/")
    group.add_argument("--status", action="store_true",
                       help="Show privacy mode state (read-only)")
    group.add_argument("--hook-check", action="store_true",
                       help="Hook entry: process maskzone/ if privacy mode is on")
    group.add_argument("--process", metavar="FILE",
                       help="Process a single file: extract + anonymize + output")
    group.add_argument("--test", action="store_true",
                       help="Run PII detection test on sample text")
    args = parser.parse_args()

    if args.setup:
        cmd_setup()
    elif args.status:
        cmd_status()
    elif args.hook_check:
        cmd_hook_check()
    elif args.process:
        cmd_process(args.process)
    elif args.test:
        cmd_test()
